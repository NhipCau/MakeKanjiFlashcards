# ===== 必要ライブラリ =====
# app.py
import pandas as pd
import streamlit as st
from deep_translator import GoogleTranslator
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pathlib import Path

#st.title("📄 漢字Flashcard作成ツール（2枚構成）") #デフォルトタイトル設定の場合はこっち

st.markdown(
    """
    <h1 style="margin-bottom:0;">
        漢字Flashcard作成ツール
        <span style="font-size:0.6em; color:gray; font-weight:normal;">（2枚構成）</span>
    </h1>
    <p style="font-size:1.2em; color:dimgray;">
        ユーザーマニュアルは 
        <a href="https://makevocabpiccard.my.canva.site/makekanjiflashcards" 
           target="_blank" 
           style="color:#1E90FF; text-decoration:none; font-weight:bold;">
           こちら
        </a>
    </p>
    """,
    unsafe_allow_html=True
)


# ── ファイル入力 ─────────────────────────────────────────
uploaded_file = st.file_uploader("Excel または CSV をアップロード", type=["xlsx", "csv"])

# ===== 設定 =====
# 列指定（番号 or 列名）
col_kanji_raw = st.text_input("漢字の列名 または番号（A列=0）", value="0")
col_hira_raw  = st.text_input("ふりがな（読み方）列名 または番号（B列=1）", value="1")

# 翻訳対象言語（カンマ区切り）
langs_str = st.text_input("翻訳対象言語（カンマ区切り）  [言語コードはこちら](https://cloud.google.com/translate/docs/languages?hl=ja)", value="en,vi,ne,my,zh-CN,zh-TW")
target_languages = [s.strip() for s in langs_str.split(",") if s.strip()]

# 位置（％指定）— 0〜100 をスライダーで
st.subheader("位置（％）設定")
kanji_y_percent = st.slider("漢字の縦位置（％）", 0, 100, 15) / 100.0
hira_y_percent  = st.slider("ふりがな（読み方）の縦位置（％）", 0, 100, 52) / 100.0
trans_y_percent = st.slider("訳語の縦位置（％）", 0, 100, 68) / 100.0

# フォントサイズ
st.subheader("フォントサイズ")
fs_kanji = st.number_input("漢字フォントサイズ", value=84, min_value=10, max_value=200)
fs_hira  = st.number_input("ふりがな（読み方）フォントサイズ", value=70, min_value=10, max_value=200)
fs_trans = st.number_input("訳語フォントサイズ", value=35, min_value=8,  max_value=120)

# スライドサイズ（EMU）
SLIDE_WIDTH  = 914400 * 10    # 10 inch
SLIDE_HEIGHT = 914400 * 7.5   # 7.5 inch

# ===== 関数 =====
def parse_col_selector(raw, df_cols):
    """数値っぽければ int、そうでなければ列名として返す"""
    try:
        i = int(raw)
        return i
    except ValueError:
        # 列名が本当に存在するか軽くチェック
        return raw

def add_textbox(slide, text, y_percent, font_size, height_percent=0.18, bold=False):
    """
    中央寄せのテキストボックスを追加。
    * ワードラップを有効化（長文でも折り返し）
    * width はスライドの 90%（左右 5% 余白）
    * height_percent はデフォ 18%（訳語エリアは複数行になりやすいので余裕）
    """
    textbox = slide.shapes.add_textbox(
        left=int(SLIDE_WIDTH * 0.02),
        top=int(SLIDE_HEIGHT * y_percent),
        width=int(SLIDE_WIDTH * 0.96),
        height=int(SLIDE_HEIGHT * height_percent),
    )
    tf = textbox.text_frame
    tf.clear()
    tf.word_wrap = True  # ← 自動折返しON

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    p.alignment = PP_ALIGN.CENTER

def add_center_line(slide):
    """スライド中央に横線（全幅）"""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        left=0,
        top=int(SLIDE_HEIGHT * 0.5),
        width=SLIDE_WIDTH,
        height=Pt(2),  # 2pt の細い帯
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0, 0, 0)
    # 枠線はいらない
    shape.line.fill.background()

def translate_word(word, lang):
    try:
        return GoogleTranslator(source="ja", target=lang).translate(word)
    except Exception:
        return f"[Error:{lang}]"

def create_ppt(df, col_kanji, col_hira, outfile_base: str):
    prs = Presentation()

    for _, row in df.iterrows():
        # 列の取得（番号 or 列名）
        kanji = str(row.iloc[col_kanji]).strip() if isinstance(col_kanji, int) else str(row[col_kanji]).strip()
        hira  = str(row.iloc[col_hira]).strip()  if isinstance(col_hira, int)  else str(row[col_hira]).strip()

        # --- Slide 1: 漢字のみ（上半分中央）
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        add_textbox(s1, kanji, kanji_y_percent, fs_kanji, height_percent=0.22, bold=True)
        add_center_line(s1)

        # --- Slide 2: 漢字 + ふりがな + 訳語
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        add_textbox(s2, kanji, kanji_y_percent, fs_kanji, height_percent=0.22, bold=True)
        add_textbox(s2, hira,  hira_y_percent,  fs_hira,  height_percent=0.20)

        translations = [translate_word(kanji, lang) for lang in target_languages]
        add_textbox(s2, "   ".join(translations), trans_y_percent, fs_trans, height_percent=0.22)

        add_center_line(s2)

    out_name = f"{outfile_base}_flashcards.pptx"
    prs.save(out_name)
    return out_name

# ── 実行 ─────────────────────────────────────────
if uploaded_file:
    # 列指定の型を決定
    col_kanji = parse_col_selector(col_kanji_raw, None)
    col_hira  = parse_col_selector(col_hira_raw, None)

    # 読み込み
    if uploaded_file.name.endswith(".xlsx"):
        df = pd.read_excel(uploaded_file)
    else:
        df = pd.read_csv(uploaded_file)

    # NaN を空文字に変換
    df = df.fillna("")

    # 出力名：アップロード元のベース名を利用
    base = Path(uploaded_file.name).stem

    if st.button("PPT を作成"):
        ppt_path = create_ppt(df, col_kanji, col_hira, base)
        with open(ppt_path, "rb") as f:
            st.download_button(
                "📥 PPTX をダウンロード",
                data=f,
                file_name=Path(ppt_path).name,
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
